import os
import logging
from PIL import Image
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
import pytesseract
from typing import Optional

# Configure logging
logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)
handler = logging.StreamHandler()
formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
handler.setFormatter(formatter)
logger.addHandler(handler)

# Configure Tesseract path
pytesseract.pytesseract.tesseract_cmd = '/opt/homebrew/bin/tesseract'

class TextExtractor:
    """Class to handle text extraction from various file formats."""
    
    def __init__(self):
        self.extractors = {
            'pdf': self._extract_from_pdf,
            'docx': self._extract_from_docx,
            'pptx': self._extract_from_pptx,
            'xlsx': self._extract_from_excel,
            'csv': self._extract_from_excel,
            'jpeg': self._extract_from_image,
            'jpg': self._extract_from_image,
            'png': self._extract_from_image
        }

    def extract(self, filepath: str) -> str:
        """
        Extract text from a file based on its extension.
        
        Args:
            filepath (str): Path to the file
            
        Returns:
            str: Extracted text or empty string if extraction fails
        """
        try:
            logger.info(f"Starting text extraction from: {filepath}")
            
            if not os.path.exists(filepath):
                raise FileNotFoundError(f"File not found: {filepath}")
                
            extension = filepath.rsplit('.', 1)[1].lower()
            extractor = self.extractors.get(extension)
            
            if not extractor:
                raise ValueError(f"No extractor available for .{extension} files")
                
            text = extractor(filepath)
            
            if not text:
                logger.warning(f"No text extracted from {filepath}")
            else:
                logger.info(f"Successfully extracted {len(text)} characters from {filepath}")
                
            return text
            
        except Exception as e:
            logger.error(f"Extraction error: {str(e)}", exc_info=True)
            return ""

    def _extract_from_pdf(self, filepath: str) -> str:
        """Extract text from PDF files."""
        logger.info(f"Extracting text from PDF: {filepath}")
        text = ""
        try:
            with open(filepath, 'rb') as file:
                reader = PdfReader(file)
                logger.info(f"PDF has {len(reader.pages)} pages")
                for i, page in enumerate(reader.pages):
                    logger.debug(f"Processing page {i+1}")
                    page_text = page.extract_text()
                    text += page_text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"PDF extraction error: {str(e)}", exc_info=True)
            raise

    def _extract_from_docx(self, filepath: str) -> str:
        """Extract text from DOCX files."""
        logger.info(f"Extracting text from DOCX: {filepath}")
        try:
            doc = Document(filepath)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            logger.info(f"Extracted {len(doc.paragraphs)} paragraphs")
            return text.strip()
        except Exception as e:
            logger.error(f"DOCX extraction error: {str(e)}", exc_info=True)
            raise

    def _extract_from_pptx(self, filepath: str) -> str:
        """Extract text from PPTX files."""
        logger.info(f"Extracting text from PPTX: {filepath}")
        try:
            prs = Presentation(filepath)
            text = []
            logger.info(f"Processing {len(prs.slides)} slides")
            for i, slide in enumerate(prs.slides):
                logger.debug(f"Processing slide {i+1}")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text).strip()
        except Exception as e:
            logger.error(f"PPTX extraction error: {str(e)}", exc_info=True)
            raise

    def _extract_from_excel(self, filepath: str) -> str:
        """Extract text from Excel files (XLSX/CSV)."""
        logger.info(f"Extracting text from Excel/CSV: {filepath}")
        try:
            if filepath.endswith('.xlsx'):
                df = pd.read_excel(filepath)
            else:
                df = pd.read_csv(filepath)
            logger.info(f"Processed {len(df)} rows and {len(df.columns)} columns")
            return df.to_string()
        except Exception as e:
            logger.error(f"Excel/CSV extraction error: {str(e)}", exc_info=True)
            raise

    def _extract_from_image(self, filepath: str) -> str:
        """Extract text from images using OCR."""
        logger.info(f"Extracting text from image using OCR: {filepath}")
        try:
            with Image.open(filepath) as img:
                logger.debug(f"Image size: {img.size}")
                text = pytesseract.image_to_string(img)
                if not text.strip():
                    logger.warning("No text detected in image")
                return text.strip()
        except Exception as e:
            logger.error(f"Image extraction error: {str(e)}", exc_info=True)
            raise